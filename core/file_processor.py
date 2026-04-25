"""
EisaX File Processor — validates and extracts text from uploaded files.

Security model:
  1. Extension allowlist — unknown extensions are rejected outright.
  2. Magic-bytes validation — file header must match declared extension.
  3. python-magic MIME check — libmagic content detection as second layer.
  4. Size limits per type before any parsing begins.
  5. ZIP-based formats (xlsx, docx, pptx) have uncompressed-size guard.
  6. Heavy parsing runs in a subprocess with a hard timeout so a
     malformed/adversarial file cannot hang the main worker.
"""

import logging
import base64
import io
import os
import traceback
import zipfile
import concurrent.futures
logger = logging.getLogger(__name__)

# ── Allowlists ────────────────────────────────────────────────────────────────

_ALLOWED_EXTENSIONS = {
    ".pdf", ".xlsx", ".xls", ".csv",
    ".pptx", ".docx",
    ".txt", ".md", ".json", ".py", ".js", ".html", ".css", ".log",
    ".yaml", ".yml", ".xml", ".sql",
    ".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff",
}

# (ext, required_magic_prefix)
_MAGIC_BYTES: dict[str, bytes] = {
    ".pdf":  b"%PDF",
    ".xlsx": b"PK\x03\x04",
    ".xls":  b"\xd0\xcf\x11\xe0",
    ".pptx": b"PK\x03\x04",
    ".docx": b"PK\x03\x04",
    ".jpg":  b"\xff\xd8\xff",
    ".jpeg": b"\xff\xd8\xff",
    ".png":  b"\x89PNG",
    ".webp": b"RIFF",
    ".bmp":  b"BM",
}

# Allowed MIME prefixes returned by python-magic, keyed on extension
_ALLOWED_MIMES: dict[str, tuple[str, ...]] = {
    ".pdf":  ("application/pdf",),
    ".xlsx": ("application/vnd.openxmlformats", "application/zip", "application/octet-stream"),
    ".xls":  ("application/vnd.ms-excel", "application/msword", "application/octet-stream", "application/x-cfb"),
    ".pptx": ("application/vnd.openxmlformats", "application/zip", "application/octet-stream"),
    ".docx": ("application/vnd.openxmlformats", "application/zip", "application/octet-stream"),
    ".jpg":  ("image/jpeg",),
    ".jpeg": ("image/jpeg",),
    ".png":  ("image/png",),
    ".webp": ("image/webp",),
    ".bmp":  ("image/bmp", "image/x-bmp"),
    ".tiff": ("image/tiff",),
    ".csv":  ("text/", "application/csv", "application/octet-stream"),
    ".txt":  ("text/",),
    ".md":   ("text/",),
    ".json": ("text/", "application/json"),
    ".xml":  ("text/", "application/xml"),
    ".html": ("text/", "application/xhtml"),
    ".sql":  ("text/", "application/x-sql"),
    ".py":   ("text/",),
    ".js":   ("text/",),
    ".css":  ("text/",),
    ".log":  ("text/",),
    ".yaml": ("text/",),
    ".yml":  ("text/",),
}

# ── Limits ────────────────────────────────────────────────────────────────────

_SIZE_LIMITS: dict[str, int] = {
    ".pdf":  10 * 1024 * 1024,
    ".xlsx":  5 * 1024 * 1024,
    ".xls":   5 * 1024 * 1024,
    ".csv":   5 * 1024 * 1024,
    ".pptx":  8 * 1024 * 1024,
    ".docx":  5 * 1024 * 1024,
    ".jpg":   8 * 1024 * 1024,
    ".jpeg":  8 * 1024 * 1024,
    ".png":   8 * 1024 * 1024,
    ".webp":  8 * 1024 * 1024,
    ".bmp":   8 * 1024 * 1024,
    ".tiff":  8 * 1024 * 1024,
}
_DEFAULT_SIZE_LIMIT = 2 * 1024 * 1024  # 2 MB for text files

_MAX_IMAGE_DIMENSION = 8000           # pixels per side
_MAX_ZIP_UNCOMPRESSED = 50 * 1024 * 1024  # 50 MB uncompressed guard
_SUBPROCESS_TIMEOUT = 20              # seconds

# ── Public entry point ────────────────────────────────────────────────────────

def process_file(filename: str, base64_data: str, content_type: str = "") -> dict:
    """
    Validate and extract text from an uploaded file.

    Returns:
        {"text": str, "type": str, "filename": str, "error": str|None}
    """
    result = {"text": "", "type": "unknown", "filename": filename, "error": None}

    try:
        raw = base64.b64decode(base64_data)
    except Exception as e:
        result["error"] = f"Base64 decode failed: {e}"
        return result

    ext = os.path.splitext(filename)[1].lower()

    # 1. Extension allowlist
    if ext not in _ALLOWED_EXTENSIONS:
        result["error"] = f"Unsupported file type: {ext or 'no extension'}"
        return result

    # 2. Size limit
    max_size = _SIZE_LIMITS.get(ext, _DEFAULT_SIZE_LIMIT)
    if len(raw) > max_size:
        result["error"] = f"File too large (max {max_size // 1024 // 1024} MB for {ext})"
        return result

    # 3. Magic-bytes validation
    required_magic = _MAGIC_BYTES.get(ext)
    if required_magic and not raw.startswith(required_magic):
        logger.warning("Magic-bytes mismatch for %s (ext=%s)", filename, ext)
        result["error"] = "File content does not match declared type"
        return result

    # 4. python-magic MIME check (best-effort; skip gracefully if unavailable)
    allowed_mimes = _ALLOWED_MIMES.get(ext)
    if allowed_mimes:
        try:
            import magic as _magic
            detected = _magic.from_buffer(raw[:4096], mime=True) or ""
            if not any(detected.startswith(m) for m in allowed_mimes):
                logger.warning("MIME mismatch for %s: detected=%s ext=%s", filename, detected, ext)
                result["error"] = f"File content (MIME: {detected}) does not match declared type"
                return result
        except ImportError:
            pass  # python-magic not installed — rely on magic-bytes above

    # 5. ZIP-bomb guard for container formats
    if ext in (".xlsx", ".docx", ".pptx") and raw.startswith(b"PK"):
        try:
            total_uncompressed = sum(
                info.file_size for info in zipfile.ZipFile(io.BytesIO(raw)).infolist()
            )
            if total_uncompressed > _MAX_ZIP_UNCOMPRESSED:
                result["error"] = "File expands too large (possible zip bomb)"
                return result
        except zipfile.BadZipFile:
            result["error"] = "Invalid ZIP/Office file"
            return result

    # 6. Parse in subprocess with timeout
    try:
        with concurrent.futures.ProcessPoolExecutor(max_workers=1) as pool:
            future = pool.submit(_parse_file, ext, raw, filename)
            text = future.result(timeout=_SUBPROCESS_TIMEOUT)
        result["text"] = text
        result["type"] = ext.lstrip(".")
    except concurrent.futures.TimeoutError:
        result["error"] = "File parsing timed out"
        logger.error("Parsing timeout for %s", filename)
    except Exception as e:
        result["error"] = f"Processing failed: {e}"
        logger.error("Parsing error for %s: %s", filename, e)

    return result


# ── Subprocess-side parser (runs in isolated process) ────────────────────────

def _parse_file(ext: str, raw: bytes, filename: str) -> str:
    """Dispatch to format reader. Executes inside a subprocess."""
    if ext == ".pdf":
        return _read_pdf(raw)
    if ext in (".xlsx", ".xls"):
        return _read_excel(raw)
    if ext == ".csv":
        return _read_csv(raw)
    if ext == ".pptx":
        return _read_pptx(raw)
    if ext == ".docx":
        return _read_docx(raw)
    if ext in (".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff"):
        return _read_image_ocr(raw, filename)
    # Plain text types
    return raw.decode("utf-8", errors="replace")


# ── Format-specific readers ───────────────────────────────────────────────────

def _read_pdf(raw: bytes) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        for i, page in enumerate(pdf.pages[:30], 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"[Page {i}]\n{page_text}")
            for table in page.extract_tables():
                if table:
                    rows = [" | ".join(str(c or "") for c in row) for row in table]
                    text_parts.append("Table:\n" + "\n".join(rows))
    return "\n\n".join(text_parts) if text_parts else "[PDF: No text extracted — may be scanned]"


def _read_excel(raw: bytes) -> str:
    import pandas as pd
    text_parts = []
    xls = pd.ExcelFile(io.BytesIO(raw))
    for sheet_name in xls.sheet_names[:10]:
        df = pd.read_excel(xls, sheet_name=sheet_name).head(5000)
        text_parts.append(f"[Sheet: {sheet_name}]\n{df.to_string(max_rows=100, max_cols=20)}")
    return "\n\n".join(text_parts)


def _read_csv(raw: bytes) -> str:
    import pandas as pd
    df = pd.read_csv(io.BytesIO(raw)).head(5000)
    return df.to_string(max_rows=200, max_cols=20)


def _read_pptx(raw: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    text_parts = []
    for i, slide in enumerate(prs.slides[:30], 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        slide_text.append(txt)
        if slide_text:
            text_parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
    return "\n\n".join(text_parts) if text_parts else "[PPTX: No text found]"


def _read_docx(raw: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(raw))
    text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        rows = [" | ".join(c.text.strip() for c in row.cells) for row in table.rows]
        if rows:
            text_parts.append("Table:\n" + "\n".join(rows))
    return "\n".join(text_parts) if text_parts else "[DOCX: No text found]"


def _read_image_ocr(raw: bytes, filename: str) -> str:
    # Dimension cap before any decoding
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(raw))
        w, h = img.size
        if w > _MAX_IMAGE_DIMENSION or h > _MAX_IMAGE_DIMENSION:
            return f"[Image: {filename} — dimensions {w}x{h} exceed limit, OCR skipped]"
    except Exception:
        pass

    # Primary: Gemini Vision
    try:
        from dotenv import load_dotenv
        load_dotenv()
        from google import genai
        from google.genai import types
        api_key = os.getenv("GEMINI_API_KEY", "")
        if api_key:
            client = genai.Client(api_key=api_key)
            ext = filename.rsplit(".", 1)[-1].lower()
            mime_map = {
                "jpg": "image/jpeg", "jpeg": "image/jpeg",
                "png": "image/png", "webp": "image/webp",
                "bmp": "image/bmp", "tiff": "image/tiff",
            }
            mime = mime_map.get(ext, "image/png")
            response = client.models.generate_content(
                model="gemini-2.0-flash",
                contents=[
                    types.Part.from_bytes(data=raw, mime_type=mime),
                    ("Extract ALL text, numbers, and table data from this image. "
                     "If it contains a portfolio or financial table, extract every "
                     "row and column precisely. Return as plain text with | separators."),
                ],
            )
            text = response.text.strip() if response.text else ""
            if text:
                return "[Img: " + filename + "]\n" + text
    except Exception:
        pass

    # Fallback: pytesseract
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(io.BytesIO(raw))
        text = pytesseract.image_to_string(img, lang="eng+ara").strip()
        if text and len(text) > 10:
            return "[OCR: " + filename + "]\n" + text
    except Exception:
        pass

    return "[Image: " + filename + " - could not extract text]"
