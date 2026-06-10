"""
File Generation Module for Clawdbot
Generates DOCX, XLSX, and PPTX files from structured data.
"""
from __future__ import annotations

import os
import re
import json
import uuid
from pathlib import Path
from typing import Any

# DOCX
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

# XLSX
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill
from openpyxl.utils import get_column_letter

# PPTX
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# PDF — optional dependency (fpdf2); falls back gracefully if not installed
try:
    from fpdf import FPDF
except ImportError:
    FPDF = object  # EISAXReport will raise NotImplementedError at runtime if used

# CSV
import csv
import io

# Export directory
EXPORT_DIR = Path(__file__).parent.parent / "static" / "exports"
EXPORT_DIR.mkdir(parents=True, exist_ok=True)

# Maximum bullets per slide
MAX_BULLETS_PER_SLIDE = 5
MAX_WORDS_PER_BULLET = 15


def render_docx_from_markdown(
    md: str,
    filename: str | None = None,
    title: str | None = None
) -> tuple[str, str]:
    """
    Convert markdown-like text to a Word document.
    
    Returns: (filename, download_url)
    """
    doc = Document()
    
    # Set default style
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    
    # Add title if provided
    if title:
        heading = doc.add_heading(title, level=0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()  # spacing
    
    # Parse markdown-like content
    for line in md.splitlines():
        stripped = line.strip()
        if not stripped:
            doc.add_paragraph()
            continue
            
        # Handle headings
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:].strip(), level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:].strip(), level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:].strip(), level=1)
        # Handle bullet points
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
        elif stripped.startswith("• "):
            doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
        # Handle numbered lists
        elif len(stripped) > 2 and stripped[0].isdigit() and stripped[1] in ".)" and stripped[2] == " ":
            doc.add_paragraph(stripped[3:].strip(), style="List Number")
        # Handle bold text markers (simple)
        elif stripped.startswith("**") and stripped.endswith("**"):
            p = doc.add_paragraph()
            run = p.add_run(stripped[2:-2])
            run.bold = True
        else:
            # Regular paragraph - handle inline bold
            p = doc.add_paragraph()
            # Simple bold handling
            if "**" in stripped:
                parts = stripped.split("**")
                for i, part in enumerate(parts):
                    if part:
                        run = p.add_run(part)
                        if i % 2 == 1:  # Odd indices are bold
                            run.bold = True
            else:
                p.add_run(stripped)
    
    # Generate filename
    safe_name = filename or f"clawdbot_{uuid.uuid4().hex[:8]}.docx"
    if not safe_name.endswith(".docx"):
        safe_name += ".docx"
    
    path = EXPORT_DIR / safe_name
    doc.save(str(path))
    
    return safe_name, f"/exports/{safe_name}"


def render_xlsx_from_table(
    table: dict,
    filename: str | None = None,
    title: str | None = None
) -> tuple[str, str]:
    """
    Create an Excel file from table data.
    
    table = {
        "sheet_name": "Data",
        "headers": ["Column1", "Column2", ...],
        "rows": [[val1, val2, ...], ...]
    }
    
    Returns: (filename, download_url)
    """
    wb = Workbook()
    ws = wb.active
    ws.title = table.get("sheet_name", "Sheet1")
    
    headers = table.get("headers", [])
    rows = table.get("rows", [])
    
    start_row = 1
    
    # Add title if provided
    if title:
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 3))
        title_cell = ws.cell(row=1, column=1, value=title)
        title_cell.font = Font(bold=True, size=14)
        title_cell.alignment = Alignment(horizontal="center")
        start_row = 3
    
    # Add headers with styling
    if headers:
        for col_idx, header in enumerate(headers, start=1):
            cell = ws.cell(row=start_row, column=col_idx, value=header)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            cell.alignment = Alignment(horizontal="center")
        start_row += 1
    
    # Add data rows
    for row_idx, row_data in enumerate(rows, start=start_row):
        for col_idx, value in enumerate(row_data, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            # Format numbers nicely
            if isinstance(value, float):
                if value < 1:  # Likely a percentage
                    cell.number_format = "0.00%"
                else:
                    cell.number_format = "#,##0.00"
    
    # Auto-size columns
    for col_idx in range(1, len(headers) + 1):
        col_letter = get_column_letter(col_idx)
        max_length = max(
            len(str(headers[col_idx - 1])) if col_idx <= len(headers) else 0,
            max((len(str(row[col_idx - 1])) if col_idx <= len(row) else 0) for row in rows) if rows else 0
        )
        ws.column_dimensions[col_letter].width = min(max(max_length + 2, 10), 50)
    
    # Generate filename
    safe_name = filename or f"clawdbot_{uuid.uuid4().hex[:8]}.xlsx"
    if not safe_name.endswith(".xlsx"):
        safe_name += ".xlsx"
    
    path = EXPORT_DIR / safe_name
    wb.save(str(path))
    
    return safe_name, f"/exports/{safe_name}"


class EISAXReport(FPDF):
    """
    Premium institutional PDF report with EISAX branding.
    """
    # Color palette
    NAVY = (10, 20, 45)
    ELECTRIC_BLUE = (0, 112, 255)
    DARK_BLUE = (20, 50, 100)
    LIGHT_GRAY = (245, 247, 250)
    MEDIUM_GRAY = (140, 150, 165)
    TEXT_DARK = (30, 40, 55)
    TEXT_BODY = (55, 65, 80)
    ACCENT_GREEN = (16, 185, 129)
    ACCENT_GOLD = (245, 180, 50)
    WHITE = (255, 255, 255)
    DIVIDER = (200, 210, 225)

    def __init__(self, report_title: str = "Report"):
        super().__init__()
        self.report_title = report_title
        self.set_auto_page_break(auto=True, margin=25)
        self._page_count_override = 0

    def header(self):
        if self.page_no() == 1:
            return  # Cover page has no header
        # Thin accent line at top
        self.set_draw_color(*self.ELECTRIC_BLUE)
        self.set_line_width(0.8)
        self.line(15, 12, self.w - 15, 12)
        # Brand name
        self.set_font("Helvetica", 'B', 8)
        self.set_text_color(*self.MEDIUM_GRAY)
        self.set_xy(15, 14)
        self.cell(0, 5, "EISAX AI", 0, 0, 'L')
        # Report title (right)
        self.set_font("Helvetica", '', 8)
        self.cell(0, 5, self.report_title.upper()[:50], 0, 0, 'R')
        self.ln(12)

    def footer(self):
        if self.page_no() == 1:
            return  # Cover page has no footer
        self.set_y(-20)
        # Divider line
        self.set_draw_color(*self.DIVIDER)
        self.set_line_width(0.3)
        self.line(15, self.h - 20, self.w - 15, self.h - 20)
        # Footer text
        self.set_font("Helvetica", '', 7)
        self.set_text_color(*self.MEDIUM_GRAY)
        from datetime import datetime
        date_str = datetime.now().strftime("%B %d, %Y")
        self.cell(0, 10, f"Generated by EISAX AI  |  {date_str}  |  Confidential", 0, 0, 'L')
        self.cell(0, 10, f"Page {self.page_no() - 1}", 0, 0, 'R')

    def _draw_cover_page(self):
        """Premium cover page with navy background and branding."""
        self.add_page()

        # Full navy background
        self.set_fill_color(*self.NAVY)
        self.rect(0, 0, self.w, self.h, 'F')

        # Electric blue accent bar (left edge)
        self.set_fill_color(*self.ELECTRIC_BLUE)
        self.rect(0, 0, 4, self.h, 'F')

        # Decorative horizontal line (upper third)
        self.set_draw_color(*self.ELECTRIC_BLUE)
        self.set_line_width(0.5)
        self.line(25, 80, self.w - 25, 80)

        # Brand identifier
        self.set_font("Helvetica", 'B', 12)
        self.set_text_color(*self.ELECTRIC_BLUE)
        self.set_xy(25, 55)
        self.cell(0, 8, "EISAX AI", 0, 1, 'L')

        self.set_font("Helvetica", '', 9)
        self.set_text_color(*self.MEDIUM_GRAY)
        self.set_xy(25, 63)
        self.cell(0, 6, "STRATEGIC INTELLIGENCE REPORT", 0, 1, 'L')

        # Main title
        self.set_font("Helvetica", 'B', 28)
        self.set_text_color(*self.WHITE)
        self.set_xy(25, 100)
        # Word-wrap the title
        self.multi_cell(self.w - 50, 14, self.report_title, 0, 'L')

        # Subtitle / date
        from datetime import datetime
        self.set_font("Helvetica", '', 11)
        self.set_text_color(*self.MEDIUM_GRAY)
        self.set_xy(25, self.get_y() + 10)
        self.cell(0, 7, datetime.now().strftime("%B %d, %Y"), 0, 1, 'L')

        # Decorative line (lower section)
        self.set_draw_color(40, 60, 100)
        self.set_line_width(0.3)
        self.line(25, self.h - 50, self.w - 25, self.h - 50)

        # Confidentiality notice
        self.set_font("Helvetica", '', 7)
        self.set_text_color(80, 100, 130)
        self.set_xy(25, self.h - 40)
        self.cell(0, 5, "CONFIDENTIAL  |  FOR AUTHORIZED RECIPIENTS ONLY", 0, 1, 'L')

    def _write_h1(self, text: str):
        """Major section heading with accent underline."""
        self.ln(6)
        self.set_font("Helvetica", 'B', 18)
        self.set_text_color(*self.TEXT_DARK)
        self.cell(0, 10, text, 0, 1, 'L')
        # Blue underline
        y = self.get_y()
        self.set_draw_color(*self.ELECTRIC_BLUE)
        self.set_line_width(1.0)
        self.line(15, y, 70, y)
        self.ln(5)

    def _write_h2(self, text: str):
        """Sub-section heading."""
        self.ln(4)
        self.set_font("Helvetica", 'B', 14)
        self.set_text_color(*self.DARK_BLUE)
        self.cell(0, 9, text, 0, 1, 'L')
        # Subtle gray underline
        y = self.get_y()
        self.set_draw_color(*self.DIVIDER)
        self.set_line_width(0.4)
        self.line(15, y, self.w - 15, y)
        self.ln(3)

    def _write_h3(self, text: str):
        """Tertiary heading."""
        self.ln(3)
        self.set_font("Helvetica", 'B', 12)
        self.set_text_color(*self.DARK_BLUE)
        self.cell(0, 8, text, 0, 1, 'L')
        self.ln(1)

    def _write_bullet(self, text: str):
        """Styled bullet point with blue marker."""
        x = self.get_x()
        # Blue bullet marker
        self.set_fill_color(*self.ELECTRIC_BLUE)
        self.ellipse(x + 4, self.get_y() + 2.5, 2.5, 2.5, 'F')
        # Bullet text
        self.set_font("Helvetica", '', 10)
        self.set_text_color(*self.TEXT_BODY)
        self.set_x(x + 12)
        self.multi_cell(self.w - 40, 6, text)
        self.ln(1)

    def _write_numbered(self, num: str, text: str):
        """Styled numbered list item."""
        x = self.get_x()
        # Number in blue circle
        self.set_font("Helvetica", 'B', 9)
        self.set_text_color(*self.ELECTRIC_BLUE)
        self.set_x(x + 3)
        self.cell(8, 6, num, 0, 0, 'C')
        # Text
        self.set_font("Helvetica", '', 10)
        self.set_text_color(*self.TEXT_BODY)
        self.set_x(x + 12)
        self.multi_cell(self.w - 40, 6, text)
        self.ln(1)

    def _write_paragraph(self, text: str):
        """Regular body paragraph."""
        self.set_font("Helvetica", '', 10)
        self.set_text_color(*self.TEXT_BODY)
        # Handle inline **bold** markers
        clean = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
        self.multi_cell(self.w - 30, 6, clean)
        self.ln(2)

    def _write_separator(self):
        """Horizontal rule."""
        self.ln(3)
        y = self.get_y()
        self.set_draw_color(*self.DIVIDER)
        self.set_line_width(0.3)
        self.line(15, y, self.w - 15, y)
        self.ln(5)

    def _write_table_row(self, cells: list, is_header: bool = False):
        """Render a markdown table row."""
        ncols = len(cells)
        if ncols == 0:
            return
        col_w = (self.w - 30) / ncols

        y_before = self.get_y()
        if is_header:
            self.set_fill_color(*self.DARK_BLUE)
            self.set_text_color(*self.WHITE)
            self.set_font("Helvetica", 'B', 9)
        else:
            # Alternate row shading
            self.set_fill_color(*self.LIGHT_GRAY)
            self.set_text_color(*self.TEXT_BODY)
            self.set_font("Helvetica", '', 9)

        for cell in cells:
            self.cell(col_w, 8, cell.strip()[:30], 1, 0, 'C', fill=is_header)
        self.ln()

    def build_from_markdown(self, md: str):
        """Parse markdown and render all content into the PDF."""
        self._draw_cover_page()
        self.add_page()

        in_table = False
        table_header_done = False
        
        for line in md.splitlines():
            stripped = line.strip()

            # Empty line
            if not stripped:
                if in_table:
                    in_table = False
                    table_header_done = False
                self.ln(3)
                continue

            # Horizontal rule
            if stripped in ("---", "***", "___"):
                in_table = False
                self._write_separator()
                continue

            # Table detection
            if "|" in stripped and stripped.startswith("|"):
                cells = [c.strip() for c in stripped.split("|") if c.strip()]
                # Skip separator rows like |---|---|
                if all(set(c) <= set("-: ") for c in cells):
                    continue
                if not in_table:
                    in_table = True
                    table_header_done = False
                    self._write_table_row(cells, is_header=True)
                    table_header_done = True
                else:
                    self._write_table_row(cells, is_header=False)
                continue
            else:
                in_table = False

            # Headings
            if stripped.startswith("### "):
                self._write_h3(stripped[4:].strip())
            elif stripped.startswith("## "):
                self._write_h2(stripped[3:].strip())
            elif stripped.startswith("# "):
                self._write_h1(stripped[2:].strip())
            # Bullet points
            elif stripped.startswith(("- ", "* ", "• ")):
                self._write_bullet(stripped[2:].strip())
            # Numbered lists
            elif len(stripped) > 2 and stripped[0].isdigit() and stripped[1] in ".)" and stripped[2] == " ":
                self._write_numbered(stripped[0], stripped[3:].strip())
            elif len(stripped) > 3 and stripped[:2].isdigit() and stripped[2] in ".)" and stripped[3] == " ":
                self._write_numbered(stripped[:2], stripped[4:].strip())
            # Bold-only line
            elif stripped.startswith("**") and stripped.endswith("**"):
                self.set_font("Helvetica", 'B', 11)
                self.set_text_color(*self.TEXT_DARK)
                self.cell(0, 8, stripped[2:-2], 0, 1)
                self.ln(1)
            # Regular paragraph
            else:
                self._write_paragraph(stripped)


def render_pdf_from_markdown(
    md: str,
    filename: str | None = None,
    title: str | None = None
) -> tuple[str, str]:
    """
    Generate a premium institutional PDF report from markdown text.
    Features: Cover page, styled headings, colored bullets, tables, branding.
    """
    report_title = title or "EISAX Report"
    pdf = EISAXReport(report_title)
    pdf.build_from_markdown(md)

    # Generate filename
    safe_name = filename or f"eisax_report_{uuid.uuid4().hex[:8]}.pdf"
    if not safe_name.endswith(".pdf"):
        safe_name += ".pdf"
    
    path = EXPORT_DIR / safe_name
    pdf.output(str(path))
    
    return safe_name, f"/exports/{safe_name}"


def render_csv_from_table(
    table: dict,
    filename: str | None = None
) -> tuple[str, str]:
    """
    Create a CSV file from table data.
    """
    headers = table.get("headers", [])
    rows = table.get("rows", [])
    
    safe_name = filename or f"clawdbot_{uuid.uuid4().hex[:8]}.csv"
    if not safe_name.endswith(".csv"):
        safe_name += ".csv"
        
    path = EXPORT_DIR / safe_name
    with open(path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        if headers:
            writer.writerow(headers)
        writer.writerows(rows)
        
    return safe_name, f"/exports/{safe_name}"


def _truncate_bullet(text: str, max_words: int = MAX_WORDS_PER_BULLET) -> str:
    """Truncate bullet text to max words for slide readability."""
    words = text.split()
    if len(words) <= max_words:
        return text
    return " ".join(words[:max_words]) + "..."


def _parse_content_to_slides(content_md: str, main_title: str) -> list[dict]:
    """
    Parse markdown content into structured slides.
    Each slide has: { "title": str, "bullets": list[str] }
    
    Rules:
    - Headings (# ## ###) become slide titles
    - Content under headings becomes bullets
    - Max 5 bullets per slide
    - Long bullets are truncated
    """
    slides = []
    current_slide = {"title": main_title, "bullets": []}
    
    lines = content_md.strip().split("\n")
    
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
            
        # Check for heading (new slide trigger)
        if stripped.startswith("### "):
            # Save current slide if it has content
            if current_slide["bullets"]:
                slides.append(current_slide)
            current_slide = {"title": stripped[4:].strip()[:60], "bullets": []}
        elif stripped.startswith("## "):
            if current_slide["bullets"]:
                slides.append(current_slide)
            current_slide = {"title": stripped[3:].strip()[:60], "bullets": []}
        elif stripped.startswith("# "):
            if current_slide["bullets"]:
                slides.append(current_slide)
            current_slide = {"title": stripped[2:].strip()[:60], "bullets": []}
        else:
            # This is content - add as bullet
            # Clean up the text
            text = stripped
            if text.startswith(("- ", "* ", "• ")):
                text = text[2:]
            elif len(text) > 2 and text[0].isdigit() and text[1] in ".)" and text[2] == " ":
                text = text[3:]
            
            # Remove markdown formatting
            text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)  # Remove bold
            text = re.sub(r'\*([^*]+)\*', r'\1', text)       # Remove italic
            text = text.strip()
            
            if text and len(text) > 3:  # Skip very short lines
                # Truncate for slide readability
                text = _truncate_bullet(text)
                
                # Check if current slide is full
                if len(current_slide["bullets"]) >= MAX_BULLETS_PER_SLIDE:
                    slides.append(current_slide)
                    current_slide = {
                        "title": f"{current_slide['title']} (cont.)",
                        "bullets": []
                    }
                
                current_slide["bullets"].append(text)
    
    # Add last slide if it has content
    if current_slide["bullets"]:
        slides.append(current_slide)
    
    # If no slides were created, make a simple one
    if not slides:
        # Split the entire text into bullets
        bullets = [_truncate_bullet(line.strip()) for line in lines if line.strip()][:MAX_BULLETS_PER_SLIDE]
        if bullets:
            slides.append({"title": main_title, "bullets": bullets})
    
    return slides


# Assets directory
ASSETS_DIR = Path(__file__).parent.parent / "static" / "assets"
PPT_BG_PATH = ASSETS_DIR / "pptx_bg.jpg"


ELECTRIC_BLUE = RGBColor(0, 112, 255)
DARK_NAVY = RGBColor(10, 20, 40)
INSIGHT_BG = RGBColor(20, 35, 60)

def render_pptx_from_slides(
    slides: list[dict],
    main_title: str = "Presentation",
    filename: str | None = None
) -> tuple[str, str]:
    """
    Create a PowerPoint presentation with premium institutional layout.
    Features: Visual Spine, Dual Layouts, and Insight Box.
    """
    prs = Presentation()
    
PPT_COVER_BG_PATH = ASSETS_DIR / "cover_bg.png"

def render_pptx_from_slides(
    slides: list[dict],
    main_title: str = "Presentation",
    filename: str | None = None
) -> tuple[str, str]:
    """
    Create a PowerPoint presentation with premium institutional layout.
    Version 1.0 Final Polish: High-res cover, increased spacing, reduced density.
    """
    prs = Presentation()
    
    # --- 1. PREMIUM COVER SLIDE ---
    title_layout = prs.slide_layouts[6] # Blank for total control
    cover_slide = prs.slides.add_slide(title_layout)
    
    if PPT_COVER_BG_PATH.exists():
        cover_slide.shapes.add_picture(str(PPT_COVER_BG_PATH), 0, 0, prs.slide_width, prs.slide_height)
    elif PPT_BG_PATH.exists():
        cover_slide.shapes.add_picture(str(PPT_BG_PATH), 0, 0, prs.slide_width, prs.slide_height)

    # Title on cover (positioned to respect the brain-circuit visual)
    # Positioning: Bottom half, slightly left-justified
    left = PptxInches(0.5)
    top = PptxInches(6.0)
    width = prs.slide_width - PptxInches(1)
    height = PptxInches(1.5)
    
    title_box = cover_slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = main_title.upper()
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.size = PptxPt(44)
    
    p2 = tf.add_paragraph()
    p2.text = "STRATEGIC ANALYSIS | INSTITUTIONAL INTELLIGENCE"
    p2.font.color.rgb = ELECTRIC_BLUE
    p2.font.size = PptxPt(14)
    p2.font.bold = True

    # --- 2. CONTENT SLIDES ---
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # 1. Apply Fixed Institutional Background
        if PPT_BG_PATH.exists():
            slide.shapes.add_picture(str(PPT_BG_PATH), 0, 0, prs.slide_width, prs.slide_height)
        
        # 2. Add Visual Spine
        spine_width = PptxInches(0.08)
        spine = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, spine_width, prs.slide_height
        )
        spine.fill.solid()
        spine.fill.fore_color.rgb = ELECTRIC_BLUE
        spine.line.fill.background()

        # 3. Add Slide Title & Subtitle
        left = PptxInches(0.6)
        top = PptxInches(0.4)
        width = prs.slide_width - PptxInches(1.2)
        height = PptxInches(1.2)
        
        # Main Title (Short & Sharp)
        title_box = slide.shapes.add_textbox(left, top, width, PptxInches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "").upper()[:60]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.size = PptxPt(32) # Slightly smaller to allow for subtitle

        # Subtitle (Contextual Info)
        subtitle_text = slide_data.get("subtitle")
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(left, top + PptxInches(0.55), width, PptxInches(0.4))
            stf = sub_box.text_frame
            sp = stf.paragraphs[0]
            sp.text = subtitle_text[:100]
            sp.font.color.rgb = ELECTRIC_BLUE
            sp.font.size = PptxPt(14)
            sp.font.bold = True
            sp.font.name = "Calibri"

        # 4. Content Layout Selection (A: Concept vs B: Tactical)
        bullets = slide_data.get("bullets", [])[:MAX_BULLETS_PER_SLIDE]
        num_bullets = len(bullets)
        
        # Reduced width for better breathing room
        margin_left = PptxInches(1.0)
        content_width = prs.slide_width - (margin_left * 2.2) 
        
        # Layout A: Concept / Overview (1-2 bullets)
        if num_bullets <= 2:
            top_content = PptxInches(2.4)
            height = PptxInches(3.3)
            
            content_box = slide.shapes.add_textbox(margin_left, top_content, content_width, height)
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, bullet_text in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet_text
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.size = PptxPt(26)
                p.font.name = "Calibri"
                p.space_after = PptxPt(30)
                p.line_spacing = 1.2
                
        # Layout B: Tactical / Lists (3-5 bullets)
        else:
            top_content = PptxInches(2.0)
            height = PptxInches(4.3)
            
            content_box = slide.shapes.add_textbox(margin_left, top_content, content_width, height)
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, bullet_text in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet_text
                p.level = 0
                p.font.color.rgb = RGBColor(220, 220, 220)
                p.font.size = PptxPt(18)
                p.font.name = "Calibri"
                p.space_after = PptxPt(16)
                p.line_spacing = 1.2

        # 5. Upgrade: Institutional Strategic Insight Call-out
        insight = slide_data.get("insight") or slide_data.get("note")
        if insight:
            box_width = prs.slide_width - PptxInches(2.5)
            box_height = PptxInches(0.6)
            left_box = (prs.slide_width - box_width) / 2
            top_box = prs.slide_height - box_height - PptxInches(0.4)
            
            # Container Pill
            rect = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left_box, top_box, box_width, box_height
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = INSIGHT_BG
            rect.line.color.rgb = ELECTRIC_BLUE
            rect.line.width = PptxPt(2.0)
            
            # Insight Text
            tf = rect.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            
            # Clean up analysis text if it was previously prefixed
            clean_insight = insight
            for prefix in ["ANALYSIS:", "AI VIEW:", "INSIGHT:", "STRATEGIC INSIGHT:"]:
                if clean_insight.upper().startswith(prefix):
                    clean_insight = clean_insight[len(prefix):].strip()
            
            p.text = f"INSTITUTIONAL INSIGHT: {clean_insight.upper()}"
            p.font.color.rgb = ELECTRIC_BLUE
            p.font.bold = True
            p.font.size = PptxPt(11)
            p.font.name = "Calibri"
    
    # Generate filename
    safe_name = filename or f"clawdbot_{uuid.uuid4().hex[:8]}.pptx"
    if not safe_name.endswith(".pptx"):
        safe_name += ".pptx"
    
    path = EXPORT_DIR / safe_name
    prs.save(str(path))
    
    return safe_name, f"/exports/{safe_name}"


def render_pptx_basic(
    title: str,
    bullets: list[str] | None = None,
    content_md: str | None = None,
    slides: list[dict] | None = None,
    filename: str | None = None
) -> tuple[str, str]:
    """
    Create a PowerPoint presentation.
    
    Priority:
    1. If slides[] is provided, use it directly (LLM-generated blueprint)
    2. If content_md is provided, parse it into slides
    3. If bullets are provided, create a simple presentation
    
    Returns: (filename, download_url)
    """
    # Priority 1: Use pre-structured slides if provided
    if slides and isinstance(slides, list) and len(slides) > 0:
        return render_pptx_from_slides(slides, title, filename)
    
    # Priority 2: Parse markdown content into slides
    if content_md:
        parsed_slides = _parse_content_to_slides(content_md, title)
        if parsed_slides:
            return render_pptx_from_slides(parsed_slides, title, filename)
    
    # Priority 3: Simple bullet-based presentation
    if bullets:
        # Split bullets into slides of max 5 each
        slide_list = []
        for i in range(0, len(bullets), MAX_BULLETS_PER_SLIDE):
            chunk = bullets[i:i+MAX_BULLETS_PER_SLIDE]
            slide_list.append({
                "title": title if i == 0 else f"{title} (cont.)",
                "bullets": [_truncate_bullet(b) for b in chunk]
            })
        return render_pptx_from_slides(slide_list, title, filename)
    
    # Fallback: Empty presentation with just title
    return render_pptx_from_slides([], title, filename)


def content_to_slide_blueprint(content_md: str, title: str = "Presentation") -> list[dict]:
    """
    Convert markdown content to a slide blueprint using smart parsing.
    This is a fallback when LLM isn't available.
    
    Returns: list of { "title": str, "bullets": list[str] }
    """
    return _parse_content_to_slides(content_md, title)


def maybe_generate_file(resp: dict) -> dict:
    """
    Check if response requests file generation, and generate if so.
    
    This is called right before returning from /chat.
    If resp.type == "file.ready", generates the file and updates resp.
    """
    if resp.get("type") != "file.ready":
        return resp
    
    data = resp.get("data") or {}
    fmt = (data.get("format") or "").lower()
    filename = data.get("filename")
    title = data.get("title", "EisaX Report")
    
    try:
        if fmt == "docx":
            content_md = data.get("content_md", "")
            name, url = render_docx_from_markdown(content_md, filename, title)
        elif fmt == "xlsx":
            table = data.get("table", {})
            name, url = render_xlsx_from_table(table, filename, title)
        elif fmt == "pptx":
            # Check for structured slides first (from LLM)
            slides = data.get("slides")
            bullets = data.get("bullets") or []
            content_md = data.get("content_md", "")
            name, url = render_pptx_basic(title, bullets, content_md, slides, filename)
        elif fmt == "pdf":
            content_md = data.get("content_md", "")
            name, url = render_pdf_from_markdown(content_md, filename, title)
        elif fmt == "csv":
            table = data.get("table", {})
            name, url = render_csv_from_table(table, filename)
        else:
            resp["reply"] = f"Unsupported file format: {fmt}. Supported: docx, xlsx, pptx, pdf, csv"
            resp["type"] = "chat.reply"
            resp["data"] = None
            return resp
        
        # Success - update response with download info
        resp["reply"] = f"Your file is ready: **{name}**"
        resp["data"] = {
            "download_url": url,
            "filename": name,
            "format": fmt
        }
        return resp
        
    except Exception as e:
        resp["type"] = "error"
        resp["reply"] = f"Failed to generate file: {str(e)}"
        resp["data"] = None
        return resp
